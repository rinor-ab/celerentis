from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from celerentis.templating import replace_tokens


def test_replace_tokens_basic(tmp_path: Path) -> None:
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tb.text_frame.text = "{{COMPANY_NAME}} / {{TAGLINE}}"
    path = tmp_path / "t.pptx"
    prs.save(str(path))  # <- pass str for mypy compatibility

    prs2 = Presentation(str(path))
    replace_tokens(prs2, {"company_name": "ACME", "tagline": "Fast", "about_bullets": []})
    text = prs2.slides[0].shapes[1].text_frame.text
    assert "ACME" in text and "Fast" in text
