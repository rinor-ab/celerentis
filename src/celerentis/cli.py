from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from rich.console import Console
from rich.table import Table
import typer

from .config import load_config
from .generator import generate
from .templating import DEFAULT_TOKENS
from .utils import iter_all_shapes

app = typer.Typer(help="Celerentis – template-driven deck generation")
console = Console()


@app.command()
def validate(config: Path):
    """Validate a YAML config file."""
    load_config(config)
    console.print("[green]✓ Config is valid[/]")


@app.command()
def inspect(template: Path):
    """List known tokens present in a template PPTX."""
    prs = Presentation(str(template))
    found: set[str] = set()
    for slide in prs.slides:
        for sh in iter_all_shapes(slide.shapes):
            if getattr(sh, "has_text_frame", False):
                txt = sh.text_frame.text or ""
                for token in DEFAULT_TOKENS:
                    if token in txt:
                        found.add(token)
    table = Table(title="Tokens in template")
    table.add_column("Token")
    for t in sorted(found):
        table.add_row(t)
    console.print(table)


@app.command(name="generate")
def generate_cmd(config: Path, output: Path | None = None):
    """Generate a deck using a YAML config."""
    cfg = load_config(config)
    if output is not None:
        cfg.output_path = output
    out = generate(cfg)
    console.print(f"[green]✓ Wrote[/] {out}")


def main():
    app()


if __name__ == "__main__":
    main()
