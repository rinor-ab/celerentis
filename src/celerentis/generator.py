from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches

from .charts import render_bar
from .config import AppConfig
from .templating import replace_tokens

EMU_PER_INCH = 914400


def _to_inches(emu: int) -> float:
    return emu / EMU_PER_INCH


def _remove_shape(shape: Any) -> None:
    elm = shape._element
    elm.getparent().remove(elm)


def _find_chart_placeholder(prs: Any, token: str) -> tuple[int, Any] | tuple[None, None]:
    for si, slide in enumerate(prs.slides):
        # flatten groups
        for shape in slide.shapes:
            stack = [shape]
            while stack:
                sh = stack.pop()
                if getattr(sh, "shape_type", None) == 6 and hasattr(sh, "shapes"):
                    stack.extend(list(sh.shapes))
                    continue
                if getattr(sh, "has_text_frame", False) and token in (sh.text_frame.text or ""):
                    return si, sh
    return None, None


def _add_logo_to_title(prs: Any, logo_path: Path) -> None:
    if not Path(logo_path).is_file():
        return
    slide = prs.slides[0]
    width = Inches(1.5)
    slide_w = int(getattr(prs, "slide_width", 0) or 0)
    left = Inches(_to_inches(slide_w) - 1.8)
    top = Inches(0.3)
    slide.shapes.add_picture(str(logo_path), left, top, width=width)


def generate(config: AppConfig) -> Path:
    prs: Any = Presentation(str(config.template_path))

    replace_tokens(
        prs,
        {
            "company_name": config.company_name,
            "tagline": config.tagline,
            "about_bullets": config.about_bullets,
        },
    )

    _add_logo_to_title(prs, config.logo_path)

    si, ph = _find_chart_placeholder(prs, config.chart.placeholder_token)
    if si is not None and ph is not None:
        chart_path = render_bar(
            config.financials,
            out_path=config.output_path.with_suffix(".chart.png"),
            title=config.chart.title,
            xlabel=config.chart.xlabel,
            ylabel=config.chart.ylabel,
        )

        slide_w = int(getattr(prs, "slide_width", 0) or 0)
        slide_h = int(getattr(prs, "slide_height", 0) or 0)
        left_emus = int(getattr(ph, "left", 0) or 0)
        top_emus = int(getattr(ph, "top", 0) or 0)
        width_emus = int(getattr(ph, "width", int(slide_w * 0.6)) or int(slide_w * 0.6))
        height_emus = int(getattr(ph, "height", int(slide_h * 0.5)) or int(slide_h * 0.5))

        left = Inches(_to_inches(left_emus))
        top = Inches(_to_inches(top_emus))
        width = Inches(_to_inches(width_emus))
        height = Inches(_to_inches(height_emus))

        _remove_shape(ph)
        prs.slides[si].shapes.add_picture(str(chart_path), left, top, width=width, height=height)

    config.output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(config.output_path))
    return config.output_path
