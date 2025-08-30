from __future__ import annotations

from dataclasses import dataclass
from typing import Dict, List, Tuple

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.text.text import _TextFrame

from .utils import iter_all_shapes

TokenValue = str | List[str]

DEFAULT_TOKENS: Dict[str, str] = {
    "{{COMPANY_NAME}}": "company_name",
    "{{TAGLINE}}": "tagline",
    "{{ABOUT_BULLETS}}": "about_bullets",
}

@dataclass(frozen=True)
class TokenStats:
    replaced: int
    missing: List[str]

def _set_bullets(tf: _TextFrame, bullets: List[str]) -> None:
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.text = str(b)
        p.level = 0

def replace_tokens(prs: Presentation, data: Dict[str, TokenValue]) -> TokenStats:
    """
    Replace occurrences of known tokens in all text frames.
    - Strings are used directly.
    - Lists are rendered as bullets.
    Returns stats for reporting.
    """
    replaced = 0
    seen: set[str] = set()

    def replace_in_shape(shape: BaseShape) -> None:
        nonlocal replaced
        if not getattr(shape, "has_text_frame", False):
            return
        tf = shape.text_frame
        if tf is None or tf.text is None:
            return
        text = tf.text
        for token, key in DEFAULT_TOKENS.items():
            if token not in text:
                continue
            seen.add(token)
            value = data.get(key, "")
            if isinstance(value, list):
                _set_bullets(tf, [str(x) for x in value])
                replaced += 1
            else:
                tf.text = text.replace(token, str(value))
                replaced += 1

    for slide in prs.slides:
        for shape in iter_all_shapes(slide.shapes):
            replace_in_shape(shape)

    missing = [t for t in DEFAULT_TOKENS if t not in seen]
    return TokenStats(replaced=replaced, missing=missing)

