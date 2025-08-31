"""PowerPoint template analyzer for extracting structure and placeholders."""

import io
from typing import List, Dict, Any, Tuple
from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Use absolute imports
from packages.core.models.slide import SlideDef, ChartToken, TemplateAnalysis


def analyze_template(pptx_bytes: bytes) -> TemplateAnalysis:
    """
    Analyze PowerPoint template to extract structure and placeholders.
    
    Args:
        pptx_bytes: Raw PowerPoint file bytes
        
    Returns:
        TemplateAnalysis with slide definitions and style mapping
    """
    try:
        print(f"Starting template analysis with {len(pptx_bytes)} bytes")
        prs = Presentation(io.BytesIO(pptx_bytes))
        print(f"Loaded presentation with {len(prs.slides)} slides")
        
        slide_defs = []
        chart_tokens = []
        style_map = {}
        
        for slide_idx, slide in enumerate(prs.slides):
            print(f"Analyzing slide {slide_idx + 1}")
            slide_def = _analyze_slide(slide, slide_idx)
            if slide_def:
                print(f"  Slide {slide_idx + 1} analysis successful: {slide_def.title}")
                slide_defs.append(slide_def)
                
                # Extract chart tokens from this slide
                slide_chart_tokens = _extract_chart_tokens(slide, slide_idx)
                chart_tokens.extend(slide_chart_tokens)
                
                # Extract style information
                slide_styles = _extract_slide_styles(slide)
                style_map.update(slide_styles)
            else:
                print(f"  Slide {slide_idx + 1} analysis failed")
        
        print(f"Template analysis completed: {len(slide_defs)} slides, {len(chart_tokens)} charts")
        
        return TemplateAnalysis(
            slide_defs=slide_defs,
            chart_tokens=chart_tokens,
            style_map=style_map
        )
        
    except Exception as e:
        print(f"Error analyzing template: {e}")
        import traceback
        traceback.print_exc()
        # Return empty analysis on error
        return TemplateAnalysis(
            slide_defs=[],
            chart_tokens=[],
            style_map={}
        )


def _analyze_slide(slide: Slide, slide_index: int) -> SlideDef:
    """Analyze individual slide for structure and content."""
    title = _extract_slide_title(slide, slide_index)
    tokens = _find_tokens(slide)
    
    return SlideDef(
        slide_index=slide_index,
        title=title,
        tokens=tokens,
        chart_tokens=[]  # Will be populated separately
    )


def _extract_slide_title(slide: Slide, slide_index: int) -> str:
    """Extract slide title from title placeholder or first text shape."""
    title = ""
    
    # Look for title placeholder
    for shape in slide.shapes:
        if isinstance(shape, SlidePlaceholder):
            try:
                # Check if placeholder_type attribute exists and is title (1)
                if hasattr(shape, 'placeholder_type') and shape.placeholder_type == 1:
                    if hasattr(shape, 'text') and shape.text:
                        title = shape.text.strip()
                        break
            except AttributeError:
                # Skip if placeholder_type doesn't exist
                continue
    
    # If no title placeholder, look for first text shape
    if not title:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                title = shape.text.strip()
                break
    
    return title or f"Slide {slide_index + 1}"


def _find_tokens(slide: Slide) -> List[str]:
    """Find text tokens (placeholders) in slide."""
    tokens = []
    
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text:
            text = shape.text.strip()
            found_tokens = _extract_text_tokens(text)
            tokens.extend(found_tokens)
    
    return list(set(tokens))  # Remove duplicates


def _extract_text_tokens(text: str) -> List[str]:
    """Extract token placeholders from text (e.g., {{COMPANY_NAME}})."""
    import re
    
    # Look for {{TOKEN}} pattern
    token_pattern = r'\{\{([^}]+)\}\}'
    matches = re.findall(token_pattern, text)
    
    return [f"{{{{{match}}}}}" for match in matches]


def _extract_chart_tokens(slide: Slide, slide_index: int) -> List[ChartToken]:
    """Extract chart tokens and their properties."""
    chart_tokens = []
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            chart_token = _analyze_chart_shape(shape, slide_index)
            if chart_token:
                chart_tokens.append(chart_token)
    
    return chart_tokens


def _analyze_chart_shape(shape: BaseShape, slide_index: int) -> ChartToken:
    """Analyze chart shape to extract token information."""
    try:
        chart = shape.chart
        
        # Try to extract chart title for token
        chart_title = ""
        if hasattr(chart, 'chart_title') and chart.chart_title:
            chart_title = chart.chart_title.text
        
        # Generate token name
        if chart_title:
            token_name = f"{{{{CHART:{chart_title}}}}}"
        else:
            token_name = f"{{{{CHART:Chart_{slide_index}}}}}"
        
        # Determine chart type
        chart_type = str(chart.chart_type) if hasattr(chart, 'chart_type') else "unknown"
        
        # Get chart dimensions
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        
        return ChartToken(
            token=token_name,
            chart_type=chart_type,
            slide_index=slide_index,
            bbox=(left, top, width, height)
        )
        
    except Exception as e:
        print(f"Error analyzing chart shape: {e}")
        return None


def _extract_slide_styles(slide: Slide) -> Dict[str, str]:
    """Extract style information from slide."""
    styles = {}
    
    # Extract background information
    if slide.background:
        bg_type = str(slide.background.type) if hasattr(slide.background, 'type') else "unknown"
        styles["background"] = bg_type
    
    # Extract layout information
    if slide.slide_layout:
        layout_name = slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "unknown"
        styles["layout"] = layout_name
    
    return styles


def _extract_text_from_shape(shape: BaseShape) -> str:
    """Extract text content from a shape."""
    if hasattr(shape, 'text'):
        return shape.text.strip()
    elif hasattr(shape, 'text_frame'):
        text_parts = []
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text:
                    text_parts.append(run.text.strip())
        return " ".join(text_parts)
    
    return ""
