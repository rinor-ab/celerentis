"""PowerPoint deck builder for creating final IM presentations."""

import io
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from PIL import Image

# Use absolute imports
from packages.core.models.slide import SlideDraft, ChartToken
from packages.core.models.financials import FinancialsData


def build_deck(
    template_bytes: bytes,
    slide_drafts: List[SlideDraft],
    financials: FinancialsData,
    logo_img_bytes: Optional[bytes] = None,
    chart_tokens: List[ChartToken] = None,
    company_name: str = "",
    website: str = ""
) -> bytes:
    """
    Build final PowerPoint deck from template and generated content.
    
    Args:
        template_bytes: Raw template PowerPoint bytes
        slide_drafts: Generated slide content
        financials: Financial data for charts
        logo_img_bytes: Company logo image bytes
        chart_tokens: Chart placeholder tokens
        company_name: Company name for replacements
        website: Company website for replacements
        
    Returns:
        Raw bytes of final PowerPoint file
    """
    try:
        # Load template
        prs = Presentation(io.BytesIO(template_bytes))
        
        # Replace text tokens
        _replace_text_tokens(prs, company_name, website)
        
        # Insert generated content
        _insert_slide_content(prs, slide_drafts)
        
        # Add charts if financial data available
        if chart_tokens and financials.series:
            _add_charts(prs, chart_tokens, financials)
        
        # Add company logo if available
        if logo_img_bytes:
            _add_company_logo(prs, logo_img_bytes)
        
        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        return output.read()
        
    except Exception as e:
        print(f"Error building deck: {e}")
        raise


def _replace_text_tokens(prs: Presentation, company_name: str, website: str):
    """Replace template tokens with actual values."""
    replacements = {
        "{{COMPANY_NAME}}": company_name,
        "{{WEBSITE}}": website,
        "{{TAGLINE}}": f"Leading {company_name} solutions",
        "{{YEAR}}": str(datetime.now().year)
    }
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                for token, value in replacements.items():
                    if token in shape.text:
                        shape.text = shape.text.replace(token, value)


def _insert_slide_content(prs: Presentation, slide_drafts: List[SlideDraft]):
    """Insert generated content into slides."""
    for slide_draft in slide_drafts:
        if slide_draft.slide_index < len(prs.slides):
            slide = prs.slides[slide_draft.slide_index]
            _update_slide_content(slide, slide_draft)


def _update_slide_content(slide: Slide, slide_draft: SlideDraft):
    """Update individual slide with generated content."""
    import re
    
    # Find text shapes and update content
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text:
            text = shape.text
            original_text = text
            
            # Handle specific template tokens with regex
            # Replace ENHANCE_BULLET_POINTS tokens
            bullet_pattern = r'\{\{ENHANCE_BULLET_POINTS:([^}]+)\}\}'
            if re.search(bullet_pattern, text):
                if slide_draft.bullet_points:
                    bullet_text = "\n• ".join([""] + slide_draft.bullet_points)
                    text = re.sub(bullet_pattern, bullet_text, text)
                else:
                    text = re.sub(bullet_pattern, "• Key points will be provided", text)
            
            # Replace ENHANCE_CONTENT tokens
            content_pattern = r'\{\{ENHANCE_CONTENT:([^}]+)\}\}'
            if re.search(content_pattern, text):
                if slide_draft.content:
                    text = re.sub(content_pattern, slide_draft.content, text)
                else:
                    text = re.sub(content_pattern, "Content will be provided", text)
            
            # Handle generic tokens as fallback
            if "{{CONTENT}}" in text:
                text = text.replace("{{CONTENT}}", slide_draft.content or "Content will be provided")
            
            if "{{BULLET_POINTS}}" in text:
                if slide_draft.bullet_points:
                    bullet_text = "\n• ".join([""] + slide_draft.bullet_points)
                    text = text.replace("{{BULLET_POINTS}}", bullet_text)
                else:
                    text = text.replace("{{BULLET_POINTS}}", "• Key points will be provided")
            
            # Handle company name and website tokens
            if "{{COMPANY_NAME}}" in text:
                text = text.replace("{{COMPANY_NAME}}", slide_draft.company_name or "Company")
            
            if "{{WEBSITE}}" in text:
                text = text.replace("{{WEBSITE}}", slide_draft.website or "www.company.com")
            
            # Only update if text actually changed
            if text != original_text:
                shape.text = text
                print(f"Updated slide {slide_draft.slide_index} shape with content")


def _add_charts(prs: Presentation, chart_tokens: List[ChartToken], financials: FinancialsData):
    """Add Excel-linked charts to slides."""
    for chart_token in chart_tokens:
        if chart_token.slide_index < len(prs.slides):
            slide = prs.slides[chart_token.slide_index]
            
            # Find chart placeholder
            chart_shape = _find_chart_placeholder(slide, chart_token)
            if chart_shape:
                _replace_chart_data(chart_shape, chart_token, financials)


def _find_chart_placeholder(slide: Slide, chart_token: ChartToken) -> Optional[BaseShape]:
    """Find chart placeholder shape in slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'chart'):
            # Check if this is the right chart by position or title
            if _is_chart_placeholder(shape, chart_token):
                return shape
    return None


def _is_chart_placeholder(shape: BaseShape, chart_token: ChartToken) -> bool:
    """Check if shape is the right chart placeholder."""
    if not hasattr(shape, 'chart'):
        return False
    
    # Check position (rough match)
    if hasattr(chart_token, 'bounding_box'):
        bbox = chart_token.bounding_box
        if bbox:
            # Allow some tolerance in position
            pos_tolerance = 50  # points
            if (abs(shape.left - bbox.get('left', 0)) < pos_tolerance and
                abs(shape.top - bbox.get('top', 0)) < pos_tolerance):
                return True
    
    # Check chart title
    if hasattr(shape.chart, 'chart_title') and shape.chart.chart_title:
        title = shape.chart.chart_title.text
        if chart_token.token in title:
            return True
    
    return False


def _replace_chart_data(chart_shape: BaseShape, chart_token: ChartToken, financials: FinancialsData):
    """Replace chart data with actual financial data."""
    try:
        chart = chart_shape.chart
        
        # Find relevant financial series
        series_name = _extract_series_name(chart_token.token)
        financial_series = _find_financial_series(financials, series_name)
        
        if financial_series:
            # Create chart data
            chart_data = ChartData()
            
            # Add categories (years)
            years = [str(year) for year, _ in financial_series.data]
            chart_data.categories = years
            
            # Add data series
            values = [value for _, value in financial_series.data]
            chart_data.add_series(series_name, values)
            
            # Apply data to chart
            chart.replace_data(chart_data)
            
    except Exception as e:
        print(f"Error updating chart data: {e}")


def _extract_series_name(token: str) -> str:
    """Extract series name from chart token."""
    if "CHART:" in token:
        return token.split("CHART:")[1].rstrip("}")
    return "Revenue"  # Default


def _find_financial_series(financials: FinancialsData, series_name: str):
    """Find financial series by name."""
    for series in financials.series:
        if series_name.lower() in series.name.lower():
            return series
    return None


def _add_company_logo(prs: Presentation, logo_bytes: bytes):
    """Add company logo to title slide."""
    try:
        # Load logo image
        logo_img = Image.open(io.BytesIO(logo_bytes))
        
        # Resize logo if needed
        max_width = 200
        max_height = 100
        
        if logo_img.width > max_width or logo_img.height > max_height:
            logo_img.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)
        
        # Save resized logo to bytes
        logo_output = io.BytesIO()
        logo_img.save(logo_output, format='PNG')
        logo_output.seek(0)
        
        # Add to first slide (title slide)
        if prs.slides:
            slide = prs.slides[0]
            
            # Find logo placeholder or add to top right
            logo_placeholder = _find_logo_placeholder(slide)
            if logo_placeholder:
                # Replace placeholder
                logo_placeholder.insert_picture(logo_output.getvalue())
            else:
                # Add logo to top right using presentation dimensions
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                # Convert from EMU to inches for positioning
                from pptx.util import Inches
                right_margin = Inches(1)  # 1 inch from right
                top_margin = Inches(0.5)  # 0.5 inch from top
                
                slide.shapes.add_picture(
                    io.BytesIO(logo_output.getvalue()),  # Convert to BytesIO object
                    slide_width - right_margin - Inches(2),  # Right side with margin
                    top_margin,                             # Top with margin
                    width=Inches(2),
                    height=Inches(1)
                )
                
    except Exception as e:
        print(f"Error adding company logo: {e}")
        import traceback
        traceback.print_exc()


def _find_logo_placeholder(slide: Slide) -> Optional[BaseShape]:
    """Find logo placeholder in slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "{{LOGO}}" in shape.text or "{{COMPANY_LOGO}}" in shape.text:
                return shape
    return None


# Import datetime for year replacement
from datetime import datetime
