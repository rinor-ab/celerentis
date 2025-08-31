"""Document bundle parser for ZIP archives."""

import io
import zipfile
from typing import List, Optional
import fitz  # PyMuPDF
import pdfplumber
from docx import Document
import xml.etree.ElementTree as ET
import traceback

# Use absolute imports
from models.document import DocChunk, DocumentBundle


def parse_bundle(zip_bytes: bytes) -> DocumentBundle:
    """
    Parse documents from ZIP archive.
    
    Args:
        zip_bytes: Raw ZIP file bytes
        
    Returns:
        DocumentBundle with extracted text chunks
    """
    chunks = []
    total_docs = 0
    total_pages = 0
    
    try:
        with zipfile.ZipFile(io.BytesIO(zip_bytes)) as zip_file:
            print(f"Processing ZIP file with {len(zip_file.filelist)} files")
            
            for file_info in zip_file.filelist:
                if file_info.is_dir():
                    continue
                    
                file_name = file_info.filename.lower()
                print(f"Processing file: {file_info.filename}")
                
                try:
                    file_bytes = zip_file.read(file_info.filename)
                    print(f"Successfully read {file_info.filename}, size: {len(file_bytes)} bytes")
                    
                    if file_name.endswith('.pdf'):
                        print(f"Parsing PDF file: {file_info.filename}")
                        doc_chunks = _parse_pdf(file_bytes, file_info.filename)
                    elif file_name.endswith('.pptx'):
                        print(f"Parsing PPTX file: {file_info.filename}")
                        doc_chunks = _parse_pptx(file_bytes, file_info.filename)
                    elif file_name.endswith('.docx'):
                        print(f"Parsing DOCX file: {file_info.filename}")
                        doc_chunks = _parse_docx(file_bytes, file_info.filename)
                    elif file_name.endswith('.txt'):
                        print(f"Parsing TXT file: {file_info.filename}")
                        doc_chunks = _parse_txt(file_bytes, file_info.filename)
                    else:
                        print(f"Skipping unsupported file type: {file_info.filename}")
                        # Skip unsupported file types
                        continue
                    
                    chunks.extend(doc_chunks)
                    total_docs += 1
                    total_pages += sum(chunk.page_count for chunk in doc_chunks)
                    print(f"Successfully processed {file_info.filename}, extracted {len(doc_chunks)} chunks")
                    
                except Exception as e:
                    print(f"Error parsing {file_info.filename}: {e}")
                    print(f"Traceback: {traceback.format_exc()}")
                    continue
            
            print(f"Bundle processing complete. Total docs: {total_docs}, Total pages: {total_pages}, Total chunks: {len(chunks)}")
            
    except Exception as e:
        print(f"Error opening ZIP file: {e}")
        print(f"Traceback: {traceback.format_exc()}")
        # Return empty bundle if ZIP processing fails
        return DocumentBundle(
            chunks=[],
            total_docs=0,
            total_pages=0
        )
    
    return DocumentBundle(
        chunks=chunks,
        total_docs=total_docs,
        total_pages=total_pages
    )


def _parse_pdf(pdf_bytes: bytes, filename: str) -> List[DocChunk]:
    """Parse PDF file using PyMuPDF and pdfplumber."""
    chunks = []
    
    print(f"Starting PDF parsing for {filename}, size: {len(pdf_bytes)} bytes")
    
    try:
        # Try PyMuPDF first (faster)
        print(f"Attempting PyMuPDF parsing for {filename}")
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        print(f"PyMuPDF successfully opened {filename}, pages: {len(doc)}")
        
        for page_num in range(len(doc)):
            try:
                page = doc.load_page(page_num)
                
                # Extract text
                text = page.get_text()
                print(f"Page {page_num + 1} text length: {len(text)} characters")
                
                # Extract images (basic)
                images = page.get_images()
                
                chunk = DocChunk(
                    text=text.strip(),
                    source_file=filename,
                    page_number=page_num + 1,
                    page_count=1,
                    chunk_type="text",
                    metadata={
                        "file_type": "pdf",
                        "page_width": page.rect.width,
                        "page_height": page.rect.height,
                        "image_count": len(images)
                    }
                )
                chunks.append(chunk)
                print(f"Successfully created chunk for page {page_num + 1}")
                
            except Exception as e:
                print(f"Error processing page {page_num + 1} in {filename}: {e}")
                continue
        
        doc.close()
        print(f"PyMuPDF parsing completed for {filename}, extracted {len(chunks)} chunks")
        
    except Exception as e:
        print(f"PyMuPDF failed for {filename}: {e}")
        print(f"PyMuPDF traceback: {traceback.format_exc()}")
        
        # Fallback to pdfplumber
        try:
            print(f"Attempting pdfplumber fallback for {filename}")
            with pdfplumber.open(io.BytesIO(pdf_bytes)) as doc:
                print(f"pdfplumber successfully opened {filename}, pages: {len(doc.pages)}")
                
                for page_num, page in enumerate(doc.pages):
                    try:
                        text = page.extract_text() or ""
                        print(f"Page {page_num + 1} text length: {len(text)} characters")
                        
                        chunk = DocChunk(
                            text=text.strip(),
                            source_file=filename,
                            page_number=page_num + 1,
                            page_count=1,
                            chunk_type="text",
                            metadata={
                                "file_type": "pdf",
                                "page_width": page.width if page.width else 0,
                                "page_height": page.height if page.height else 0
                            }
                        )
                        chunks.append(chunk)
                        print(f"Successfully created chunk for page {page_num + 1} using pdfplumber")
                        
                    except Exception as e:
                        print(f"Error processing page {page_num + 1} with pdfplumber in {filename}: {e}")
                        continue
                        
            print(f"pdfplumber parsing completed for {filename}, extracted {len(chunks)} chunks")
            
        except Exception as e:
            print(f"pdfplumber also failed for {filename}: {e}")
            print(f"pdfplumber traceback: {traceback.format_exc()}")
    
    print(f"Final result for {filename}: {len(chunks)} chunks extracted")
    return chunks


def _parse_pptx(pptx_bytes: bytes, filename: str) -> List[DocChunk]:
    """Parse PowerPoint file."""
    chunks = []
    
    try:
        from pptx import Presentation
        
        prs = Presentation(io.BytesIO(pptx_bytes))
        
        for slide_num, slide in enumerate(prs.slides):
            text_parts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text.strip())
            
            slide_text = " ".join(text_parts)
            
            if slide_text.strip():
                chunk = DocChunk(
                    text=slide_text,
                    source_file=filename,
                    page_number=slide_num + 1,
                    page_count=1,
                    chunk_type="slide",
                    metadata={
                        "file_type": "pptx",
                        "slide_layout": slide.slide_layout.name if slide.slide_layout else "unknown"
                    }
                )
                chunks.append(chunk)
        
    except Exception as e:
        print(f"Failed to parse PPTX {filename}: {e}")
    
    return chunks


def _parse_docx(docx_bytes: bytes, filename: str) -> List[DocChunk]:
    """Parse Word document."""
    chunks = []
    
    try:
        doc = Document(io.BytesIO(docx_bytes))
        
        # Extract text from paragraphs
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())
        
        if text_parts:
            chunk = DocChunk(
                text="\n".join(text_parts),
                source_file=filename,
                page_number=1,
                page_count=1,
                chunk_type="document",
                metadata={
                    "file_type": "docx",
                    "paragraph_count": len(text_parts)
                }
            )
            chunks.append(chunk)
        
    except Exception:
        # Fallback to XML parsing
        try:
            chunks = _parse_docx_xml(docx_bytes, filename)
        except Exception as e:
            print(f"Failed to parse DOCX {filename}: {e}")
    
    return chunks


def _parse_docx_xml(docx_bytes: bytes, filename: str) -> List[DocChunk]:
    """Parse Word document using XML fallback."""
    chunks = []
    
    try:
        with zipfile.ZipFile(io.BytesIO(docx_bytes)) as zip_file:
            # Extract word/document.xml
            xml_content = zip_file.read("word/document.xml")
            root = ET.fromstring(xml_content)
            
            # Extract text from w:t elements
            text_parts = []
            for text_elem in root.findall(".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"):
                if text_elem.text:
                    text_parts.append(text_elem.text.strip())
            
            if text_parts:
                chunk = DocChunk(
                    text=" ".join(text_parts),
                    source_file=filename,
                    page_number=1,
                    page_count=1,
                    chunk_type="document",
                    metadata={
                        "file_type": "docx",
                        "parsing_method": "xml_fallback"
                    }
                )
                chunks.append(chunk)
        
    except Exception as e:
        print(f"Failed to parse DOCX XML {filename}: {e}")
    
    return chunks


def _parse_txt(txt_bytes: bytes, filename: str) -> List[DocChunk]:
    """Parse text file."""
    try:
        text = txt_bytes.decode('utf-8', errors='ignore')
        
        chunk = DocChunk(
            text=text.strip(),
            source_file=filename,
            page_number=1,
            page_count=1,
            chunk_type="text",
            metadata={
                "file_type": "txt",
                "encoding": "utf-8"
            }
        )
        
        return [chunk]
        
    except Exception as e:
        print(f"Failed to parse TXT {filename}: {e}")
        return []
