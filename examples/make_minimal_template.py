import os

# Ensure headless backend BEFORE importing pyplot
os.environ["MPLBACKEND"] = "Agg"

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches


prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
t1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
t1.text_frame.text = "{{COMPANY_NAME}}"
t2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
t2.text_frame.text = "{{TAGLINE}}"

# About slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(0.5)).text_frame.text = "About"
body = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(3))
body.text_frame.text = "{{ABOUT_BULLETS}}"

# Financials slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(0.5)).text_frame.text = "Financials"
ph = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
ph.text_frame.text = "{{CHART_PLACEHOLDER}}"

prs.save("examples/template.pptx")

# simple square logo
plt.figure(figsize=(1, 1))
plt.text(0.5, 0.5, "LOGO", ha="center", va="center")
plt.axis("off")
plt.tight_layout()
plt.savefig("examples/logo.png")
plt.close()

print("Wrote examples/template.pptx and examples/logo.png")
